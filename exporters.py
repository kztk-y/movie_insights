"""
Movie Insights - Export Functions
Excel・PowerPoint出力機能
"""

import io
import zipfile
from pathlib import Path
from typing import List, Tuple

from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
from openpyxl.utils import get_column_letter
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from scene_detector import SceneInfo


def export_to_excel(
    scenes: List[SceneInfo],
    video_info: dict,
    output_path: str,
    thumbnail_size: Tuple[int, int] = (160, 90)
) -> str:
    """
    シーン一覧をExcelファイルに出力

    Args:
        scenes: シーン情報のリスト
        video_info: 動画の基本情報
        output_path: 出力ファイルパス
        thumbnail_size: サムネイルサイズ (width, height)

    Returns:
        出力ファイルパス
    """
    wb = Workbook()
    ws = wb.active
    ws.title = "シーン一覧"

    # スタイル定義
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin")
    )
    center_align = Alignment(horizontal="center", vertical="center")

    # ヘッダー行
    headers = ["No.", "サムネイル", "開始時間", "終了時間", "長さ(秒)", "メモ"]
    col_widths = [6, 25, 15, 15, 12, 30]

    for col, (header, width) in enumerate(zip(headers, col_widths), 1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.border = border
        cell.alignment = center_align
        ws.column_dimensions[get_column_letter(col)].width = width

    # データ行
    row_height = 70  # サムネイル用の行高さ

    for i, scene in enumerate(scenes, 2):
        # 行の高さを設定
        ws.row_dimensions[i].height = row_height

        # No.
        ws.cell(row=i, column=1, value=scene.scene_num).alignment = center_align
        ws.cell(row=i, column=1).border = border

        # サムネイル
        ws.cell(row=i, column=2).border = border
        if scene.thumbnail_path and Path(scene.thumbnail_path).exists():
            # サムネイルをリサイズして挿入
            img = Image.open(scene.thumbnail_path)
            img.thumbnail(thumbnail_size, Image.Resampling.LANCZOS)

            # 一時ファイルに保存
            img_buffer = io.BytesIO()
            img.save(img_buffer, format="PNG")
            img_buffer.seek(0)

            xl_img = XLImage(img_buffer)
            ws.add_image(xl_img, f"B{i}")

        # 開始時間
        ws.cell(row=i, column=3, value=scene.start_timecode).alignment = center_align
        ws.cell(row=i, column=3).border = border

        # 終了時間
        ws.cell(row=i, column=4, value=scene.end_timecode).alignment = center_align
        ws.cell(row=i, column=4).border = border

        # 長さ
        ws.cell(row=i, column=5, value=round(scene.duration, 2)).alignment = center_align
        ws.cell(row=i, column=5).border = border

        # メモ（空欄）
        ws.cell(row=i, column=6).border = border

    # 動画情報シートを追加
    info_ws = wb.create_sheet(title="動画情報")
    info_data = [
        ("ファイル名", Path(video_info["path"]).name if video_info["path"] else ""),
        ("総再生時間", video_info.get("duration_formatted", "")),
        ("FPS", video_info.get("fps", "")),
        ("総フレーム数", video_info.get("total_frames", "")),
        ("検出シーン数", video_info.get("scene_count", ""))
    ]
    for row, (label, value) in enumerate(info_data, 1):
        info_ws.cell(row=row, column=1, value=label).font = Font(bold=True)
        info_ws.cell(row=row, column=2, value=value)

    info_ws.column_dimensions["A"].width = 15
    info_ws.column_dimensions["B"].width = 30

    wb.save(output_path)
    return output_path


def export_to_pptx(
    scenes: List[SceneInfo],
    video_info: dict,
    output_path: str,
    images_per_slide: int = 6,
    grid_cols: int = 3
) -> str:
    """
    シーンをPowerPointスライドに出力

    Args:
        scenes: シーン情報のリスト
        video_info: 動画の基本情報
        output_path: 出力ファイルパス
        images_per_slide: 1スライドあたりの画像数
        grid_cols: グリッドの列数

    Returns:
        出力ファイルパス
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # タイトルスライド
    title_layout = prs.slide_layouts[6]  # 空白レイアウト
    title_slide = prs.slides.add_slide(title_layout)

    # タイトルテキスト
    title_box = title_slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Scene Analysis Report"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # サブタイトル
    subtitle_box = title_slide.shapes.add_textbox(
        Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    video_name = Path(video_info["path"]).name if video_info["path"] else "Unknown"
    subtitle_para.text = f"{video_name}"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 情報テキスト
    info_box = title_slide.shapes.add_textbox(
        Inches(0.5), Inches(5), Inches(12.333), Inches(1)
    )
    info_frame = info_box.text_frame
    info_para = info_frame.paragraphs[0]
    info_para.text = f"Duration: {video_info.get('duration_formatted', '')} | Scenes: {video_info.get('scene_count', 0)}"
    info_para.font.size = Pt(18)
    info_para.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    info_para.alignment = PP_ALIGN.CENTER

    # シーングリッドスライド
    grid_rows = (images_per_slide + grid_cols - 1) // grid_cols

    # スライドのマージンとサイズ計算
    margin_x = 0.3
    margin_y = 0.5
    gap = 0.15
    available_width = 13.333 - (margin_x * 2) - (gap * (grid_cols - 1))
    available_height = 7.5 - (margin_y * 2) - (gap * (grid_rows - 1)) - 0.5  # タイトル用スペース

    img_width = available_width / grid_cols
    img_height = available_height / grid_rows

    # 16:9アスペクト比を維持
    if img_width / img_height > 16 / 9:
        img_width = img_height * 16 / 9
    else:
        img_height = img_width * 9 / 16

    # シーンをグループ化してスライド作成
    for slide_num, start_idx in enumerate(range(0, len(scenes), images_per_slide)):
        slide_scenes = scenes[start_idx:start_idx + images_per_slide]

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # スライドタイトル
        slide_title = slide.shapes.add_textbox(
            Inches(margin_x), Inches(0.2), Inches(12), Inches(0.4)
        )
        slide_title_frame = slide_title.text_frame
        slide_title_para = slide_title_frame.paragraphs[0]
        slide_title_para.text = f"Scenes {start_idx + 1} - {start_idx + len(slide_scenes)}"
        slide_title_para.font.size = Pt(18)
        slide_title_para.font.bold = True

        # 画像をグリッド配置
        for i, scene in enumerate(slide_scenes):
            row = i // grid_cols
            col = i % grid_cols

            x = margin_x + col * (img_width + gap)
            y = margin_y + 0.3 + row * (img_height + gap + 0.3)  # ラベル用スペース

            if scene.thumbnail_path and Path(scene.thumbnail_path).exists():
                # 画像を追加
                pic = slide.shapes.add_picture(
                    scene.thumbnail_path,
                    Inches(x),
                    Inches(y),
                    Inches(img_width),
                    Inches(img_height)
                )

                # シーン番号ラベル
                label_box = slide.shapes.add_textbox(
                    Inches(x), Inches(y + img_height + 0.02),
                    Inches(img_width), Inches(0.25)
                )
                label_frame = label_box.text_frame
                label_para = label_frame.paragraphs[0]
                label_para.text = f"#{scene.scene_num} | {scene.start_timecode} ({scene.duration:.1f}s)"
                label_para.font.size = Pt(10)
                label_para.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return output_path


def export_images_zip(
    scenes: List[SceneInfo],
    output_path: str
) -> str:
    """
    サムネイル画像をZIPファイルに圧縮

    Args:
        scenes: シーン情報のリスト
        output_path: 出力ファイルパス

    Returns:
        出力ファイルパス
    """
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for scene in scenes:
            if scene.thumbnail_path and Path(scene.thumbnail_path).exists():
                arcname = Path(scene.thumbnail_path).name
                zf.write(scene.thumbnail_path, arcname)

    return output_path
